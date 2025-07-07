
import os
import tempfile
import logging
from pathlib import Path
from typing import List, Dict, Set
from dataclasses import dataclass
from pptx import Presentation
from pydrive2.auth import GoogleAuth
from pydrive2.drive import GoogleDrive
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.documents import Document
from dotenv import load_dotenv
from PDFProcessor import PDFProcessor
from llama_index.core.node_parser import SemanticSplitterNodeParser
from llama_index.embeddings.langchain import LangchainEmbedding
from llama_index.core import Document as LlamaDocument

load_dotenv(override=True)
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

@dataclass
class ProcessedFile:
    filename: str
    status: str
    chunks: int = 0
    message: str = ""
    file_type: str = ""

class DriveIngestionError(Exception):
    """Custom exception for drive ingestion errors"""
    pass

def init_drive() -> GoogleDrive:
    """Initialize Google Drive API client with service account"""
    try:
        service_account_file = os.getenv("GOOGLE_SERVICE_ACCOUNT_FILE")
        if not service_account_file:
            raise DriveIngestionError("GOOGLE_SERVICE_ACCOUNT_FILE environment variable not set")
        
        if not os.path.exists(service_account_file):
            raise DriveIngestionError(f"Service account file not found: {service_account_file}")
        
        settings = {
            "client_config_backend": "service",
            "service_config": {
                "client_json_file_path": service_account_file
            }
        }
        
        gauth = GoogleAuth(settings=settings)
        gauth.ServiceAuth()
        return GoogleDrive(gauth)
    except Exception as e:
        logger.error(f"Google Drive initialization failed: {str(e)}")
        raise DriveIngestionError(f"Drive authentication failed: {str(e)}")

def load_pptx(file_path: str) -> List[Document]:
    """Extract text from PowerPoint presentation"""
    try:
        prs = Presentation(file_path)
        text = "\n".join(
            shape.text for slide in prs.slides 
            for shape in slide.shapes 
            if hasattr(shape, "text")
        )
        return [Document(page_content=text, metadata={"source": file_path})]
    except Exception as e:
        logger.error(f"Failed to process PPTX file {file_path}: {str(e)}")
        raise DriveIngestionError(f"PPTX processing failed: {str(e)}")

def split_documents(documents: List[Dict], chunk_size: int = 300) -> List:
    """Split documents using semantic chunking"""
    try:

        llama_docs = [
            LlamaDocument(text=doc.page_content, metadata=doc.metadata)
            for doc in documents
        ]

        embedding = get_embedding_model()
        langchain_embed_wrapper = LangchainEmbedding(embedding)

        splitter = SemanticSplitterNodeParser(
            embed_model=langchain_embed_wrapper,
            chunk_size=chunk_size  # token count
        )

        nodes = []
        for doc in llama_docs:
            nodes.extend(splitter.get_nodes_from_documents([doc]))

        from langchain_core.documents import Document as LangchainDoc
        return [
            LangchainDoc(page_content=node.text, metadata=node.metadata)
            for node in nodes
        ]
    except Exception as e:
        logger.error(f"Semantic document splitting failed: {str(e)}")
        raise DriveIngestionError(f"Semantic splitting failed: {str(e)}")


def get_embedding_model() -> HuggingFaceEmbeddings:
    """Initialize the embedding model"""
    try:
        return HuggingFaceEmbeddings(model_name='all-MiniLM-L6-v2')
    except Exception as e:
        logger.error(f"Failed to initialize embedding model: {str(e)}")
        raise DriveIngestionError("Embedding model initialization failed")

def store_embeddings(docs: List, embedding_model: HuggingFaceEmbeddings, persist_directory: str) -> None:
    """Store document embeddings in ChromaDB"""
    try:
        db = Chroma.from_documents(
            documents=docs,
            embedding=embedding_model,
            persist_directory=persist_directory
        )
        db.persist()
    except Exception as e:
        logger.error(f"Failed to store embeddings: {str(e)}")
        raise DriveIngestionError(f"Embedding storage failed: {str(e)}")

def load_processed_set(tracker_txt: str) -> Set[str]:
    """Load set of already processed files"""
    try:
        if not os.path.exists(tracker_txt):
            with open(tracker_txt, 'w') as f:
                pass  
            return set()
        
        with open(tracker_txt, 'r') as f:
            return set(line.strip() for line in f if line.strip())
    except Exception as e:
        logger.error(f"Failed to load processed files tracker: {str(e)}")
        raise DriveIngestionError("Processed files tracker loading failed")

def mark_as_processed(filename: str, tracker_txt: str) -> None:
    """Mark a file as processed in the tracker"""
    try:
        with open(tracker_txt, 'a') as f:
            f.write(f"{filename}\n")
    except Exception as e:
        logger.error(f"Failed to update processed files tracker: {str(e)}")
        raise DriveIngestionError("Processed files tracker update failed")



def process_drive_folder(
    folder_id: str,
    persist_directory: str,
    tracker_txt: str,
    force_reprocess: bool = False
) -> List[ProcessedFile]:
    try:
        logger.info(f"Initializing Drive with folder ID: {folder_id}")
        drive = init_drive()

        logger.info("Listing files from Google Drive folder")
        file_list = drive.ListFile({'q': f"'{folder_id}' in parents and trashed=false"}).GetList()
        logger.info(f"Total files fetched from Drive: {len(file_list)}")

        if not os.path.exists(tracker_txt):
            logger.info(f"KB Agent Tracker file not found. Creating: {tracker_txt}")
            Path(tracker_txt).touch()

        processed_files = load_processed_set(tracker_txt)
        logger.info(f"Already processed files loaded: {len(processed_files)} entries")

        embedding_model = get_embedding_model()
        logger.info("Embedding model initialized")

        results = []

        for file in file_list:
            file_name = file['title']
            file_ext = Path(file_name).suffix.lower()
            logger.info(f"Evaluating file: {file_name} (ext: {file_ext})")

            if file_ext not in ('.pdf', '.pptx'):
                logger.info(f"Skipping unsupported file: {file_name}")
                continue

            file_type = 'pdf' if file_ext == '.pdf' else 'pptx'
            result = ProcessedFile(filename=file_name, status='skipped', file_type=file_type)

            if not force_reprocess and file_name in processed_files:
                logger.info(f"Skipping already processed file: {file_name}")
                results.append(result)
                continue

            with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp:
                try:
                    logger.info(f"Downloading file: {file_name}")
                    file.GetContentFile(tmp.name)
                    temp_path = tmp.name

                    logger.info(f"Parsing document: {file_name}")
                    if file_ext == '.pdf':
                        try:
                            processor = PDFProcessor(
                                pdf_path=temp_path,
                                output_dir=os.path.dirname(temp_path),
                                output_filename=f"{Path(temp_path).stem}_output.json"
                            )
                            final_data = processor.process(cleanup_temp=True)

                            # Save output.md manually
                            md_output_path = os.path.join(os.path.dirname(temp_path), f"{Path(temp_path).stem}_output.md")
                            with open(md_output_path, "w", encoding="utf-8") as f_md:
                                for page in final_data:
                                    f_md.write(page["text"].strip() + "\n\n")

                            # Convert to Langchain documents
                            chunks = [
                                Document(page_content=page["text"], metadata={
                                    "source": file_name,
                                    "page": page.get("metadata", {}).get("page"),
                                    "section": page.get("metadata", {}).get("main_section_header"),
                                })
                                for page in final_data
                            ]

                        except Exception as e:
                            logger.warning(f"PDFProcessor failed for {file_name}, falling back to PyPDFLoader: {str(e)}")
                            docs = PyPDFLoader(temp_path).load()
                            for doc in docs:
                                doc.metadata["source"] = file_name
                            chunks = split_documents(docs)
                    elif file_ext == '.pptx':
                        docs = load_pptx(temp_path)

                    for doc in docs:
                        doc.metadata["source"] = file_name

                    logger.info(f"Splitting document: {file_name}")
                    chunks = split_documents(docs)

                    logger.info(f"Storing embeddings for: {file_name}")
                    store_embeddings(chunks, embedding_model, persist_directory)

                    if not force_reprocess:
                        mark_as_processed(file_name, tracker_txt)
                        logger.info(f"Marked as processed: {file_name}")

                    result.status = 'processed'
                    result.chunks = len(chunks)
                    results.append(result)
                    logger.info(f"Successfully processed: {file_name} (chunks: {len(chunks)})")

                except Exception as e:
                    result.status = 'error'
                    result.message = str(e)
                    results.append(result)
                    logger.error(f"Failed to process {file_name}: {str(e)}")

                finally:
                    try:
                        os.remove(temp_path)
                        logger.info(f"Deleted temp file: {temp_path}")
                    except Exception as e:
                        logger.warning(f"Failed to delete temp file {temp_path}: {str(e)}")

        logger.info(f"Processing completed. Total results: {len(results)}")
        return results

    except Exception as e:
        logger.error(f"Drive folder processing failed: {str(e)}")
        raise DriveIngestionError(f"Drive ingestion failed: {str(e)}")
